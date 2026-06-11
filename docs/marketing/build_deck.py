"""Builds RocketRFP customer pitch deck (.pptx).

Brand:
  Deep Space Navy   #0B1B3B
  Ignition Orange   #FF6B35
  Launch Blue       #3B82F6
  Bg Light          #F8FAFC
  Body Text         #1E293B
  Muted             #64748B
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

NAVY = RGBColor(0x0B, 0x1B, 0x3B)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
BG = RGBColor(0xF8, 0xFA, 0xFC)
BODY = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HERE = Path(__file__).parent
DIAG = HERE / "diagrams"
OUT = HERE / "RocketRFP-Pitch-Deck.pptx"

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=BODY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Inter"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w is not None:
            s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_round(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.12
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def header_bar(slide, title, kicker=None):
    # left orange accent
    add_rect(slide, Inches(0.5), Inches(0.5), Inches(0.12), Inches(0.6), fill=ORANGE)
    if kicker:
        add_text(slide, Inches(0.75), Inches(0.45), Inches(10), Inches(0.3),
                 kicker.upper(), size=11, bold=True, color=ORANGE)
        add_text(slide, Inches(0.75), Inches(0.78), Inches(12), Inches(0.5),
                 title, size=28, bold=True, color=NAVY)
    else:
        add_text(slide, Inches(0.75), Inches(0.55), Inches(12), Inches(0.6),
                 title, size=30, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, num, total):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(4), Inches(0.3),
             "RocketRFP", size=10, bold=True, color=NAVY)
    add_text(slide, Inches(8.83), Inches(7.05), Inches(4), Inches(0.3),
             f"Confidential  ·  {num} / {total}", size=10, color=MUTED,
             align=PP_ALIGN.RIGHT)


def wordmark(slide, x, y, size=44, color=NAVY, accent=ORANGE):
    """Text wordmark: 'Rocket' in navy, 'RFP' in orange, with chevron."""
    tb = slide.shapes.add_textbox(x, y, Inches(6), Inches(1.2))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = "Rocket"
    r1.font.name = "Inter"; r1.font.size = Pt(size); r1.font.bold = True
    r1.font.color.rgb = color
    r2 = p.add_run(); r2.text = "RFP"
    r2.font.name = "Inter"; r2.font.size = Pt(size); r2.font.bold = True
    r2.font.color.rgb = accent
    return tb


# ============================================================
# SLIDE 1 — Title
# ============================================================
def slide_title(num, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, NAVY)
    # Orange diagonal accent strip
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7))
    accent.fill.solid(); accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background(); accent.shadow.inherit = False
    # Wordmark — light variant
    wordmark(s, Inches(1.0), Inches(2.4), size=72, color=WHITE, accent=ORANGE)
    add_text(s, Inches(1.0), Inches(3.6), Inches(11), Inches(0.6),
             "AI-Powered Proposal Response, From Days to Hours",
             size=24, color=WHITE, font="Inter")
    add_text(s, Inches(1.0), Inches(4.4), Inches(11), Inches(0.5),
             "Win more bids without growing the team.",
             size=18, color=BG, font="Inter")
    add_text(s, Inches(1.0), Inches(6.95), Inches(11), Inches(0.4),
             "Customer Briefing", size=12, bold=True, color=NAVY,
             align=PP_ALIGN.LEFT)
    return s


# ============================================================
# SLIDE 2 — The Monday Morning Problem
# ============================================================
def slide_problem(num, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    header_bar(s, "Monday morning, your proposal team is buried.", kicker="The problem")

    pains = [
        ("⏱️", "Days lost per response",
         "A single RFP eats 3–7 days of senior time, copy-pasted from old win files."),
        ("📄", "Format chaos",
         "PDFs, Word docs, web portals with character limits — no single tool handles them all."),
        ("🔍", "Knowledge scattered",
         "Past projects, certs, bios, and pricing live in Slack, email, and someone's desktop."),
        ("📉", "Bids left on the table",
         "Most teams respond to under 20% of qualified opportunities — pure capacity loss."),
    ]
    x0, y0 = Inches(0.6), Inches(1.7)
    cw, ch = Inches(6.05), Inches(2.55)
    gap = Inches(0.15)
    for i, (icon, title, body) in enumerate(pains):
        col = i % 2; row = i // 2
        x = x0 + col * (cw + gap)
        y = y0 + row * (ch + gap)
        card = add_round(s, x, y, cw, ch, fill=WHITE, line=NAVY)
        card.line.width = Pt(0.75)
        add_text(s, x + Inches(0.3), y + Inches(0.25), Inches(1), Inches(0.6),
                 icon, size=32)
        add_text(s, x + Inches(1.2), y + Inches(0.3), cw - Inches(1.4), Inches(0.5),
                 title, size=20, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), y + Inches(1.15), cw - Inches(0.6), ch - Inches(1.3),
                 body, size=15, color=BODY)
    footer(s, num, total)
    return s


# ============================================================
# SLIDE 3 — The Cost
# ============================================================
def slide_cost(num, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    header_bar(s, "What it actually costs you.", kicker="The math")

    metrics = [
        ("$5K–$25K", "Labor cost per RFP response", BLUE),
        ("3–7 days", "Senior time per response", ORANGE),
        ("<20%", "Of qualified bids your team can answer", NAVY),
    ]
    x0, y = Inches(0.6), Inches(1.9)
    cw, ch = Inches(4.05), Inches(2.6)
    gap = Inches(0.15)
    for i, (big, label, color) in enumerate(metrics):
        x = x0 + i * (cw + gap)
        card = add_round(s, x, y, cw, ch, fill=WHITE, line=color)
        card.line.width = Pt(2)
        add_text(s, x, y + Inches(0.5), cw, Inches(1.1),
                 big, size=54, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), y + Inches(1.75), cw - Inches(0.6), Inches(0.8),
                 label, size=15, color=BODY, align=PP_ALIGN.CENTER)

    add_round(s, Inches(0.6), Inches(4.85), Inches(12.13), Inches(1.7),
              fill=NAVY, line=ORANGE)
    add_text(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.5),
             "If your team could answer 3× more RFPs with the same headcount,",
             size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.9), Inches(5.55), Inches(11.5), Inches(0.7),
             "what would that do to your pipeline?",
             size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    footer(s, num, total)
    return s


# ============================================================
# SLIDE 4 — How It Works (Dual Loop diagram)
# ============================================================
def slide_dual_loop(num, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    header_bar(s, "Two loops. One smarter proposal every time.", kicker="How it works")
    s.shapes.add_picture(str(DIAG / "01-dual-loop.png"),
                         Inches(0.6), Inches(1.7),
                         width=Inches(12.13), height=Inches(4.4))
    add_text(s, Inches(0.6), Inches(6.25), Inches(12.13), Inches(0.6),
             "Your services are standardized. Customer RFPs are variable. RocketRFP is the intelligent mapping in between.",
             size=14, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, num, total)
    return s


# ============================================================
# SLIDE 5 — Loop 1 Smart Intake
# ============================================================
def slide_intake(num, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    header_bar(s, "Loop 1 — Smart Intake", kicker="Turn 50 pages into a checklist")

    # Left: before
    add_round(s, Inches(0.6), Inches(1.85), Inches(5.8), Inches(4.6),
              fill=WHITE, line=MUTED)
    add_text(s, Inches(0.85), Inches(2.0), Inches(5.3), Inches(0.4),
             "BEFORE", size=11, bold=True, color=MUTED)
    add_text(s, Inches(0.85), Inches(2.4), Inches(5.3), Inches(3.9),
             "📄  A 50-page PDF.\n\n"
             "🌐  A web portal with questions buried in tabs.\n\n"
             "📋  Hours of reading just to know what's being asked.\n\n"
             "❓  Easy to miss a required cert, deliverable, or page-limit rule.",
             size=15, color=BODY)

    # Right: after
    after = add_round(s, Inches(6.95), Inches(1.85), Inches(5.8), Inches(4.6),
                      fill=NAVY, line=ORANGE)
    after.line.width = Pt(2)
    add_text(s, Inches(7.2), Inches(2.0), Inches(5.3), Inches(0.4),
             "AFTER  ·  ~2 MIN", size=11, bold=True, color=ORANGE)
    add_text(s, Inches(7.2), Inches(2.4), Inches(5.3), Inches(3.9),
             "✅  27 questions extracted and tagged\n\n"
             "✅  6 required deliverables flagged\n\n"
             "✅  Submission rules captured (page limits, format, due date)\n\n"
             "✅  Human reviews the RFP profile before any drafting starts",
             size=15, color=WHITE)
    footer(s, num, total)
    return s


# ============================================================
# SLIDE 6 — Loop 2 Drafting Canvas
# ============================================================
def slide_canvas(num, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    header_bar(s, "Loop 2 — The Drafting Canvas", kicker="Where the work actually happens")

    # Mock canvas: left = question, right = AI draft
    add_text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.4),
             "Side-by-side editor. AI proposes. You approve, edit, or chat with your knowledge base.",
             size=14, color=MUTED)

    # Left pane
    add_round(s, Inches(0.6), Inches(2.2), Inches(5.9), Inches(3.6),
              fill=WHITE, line=NAVY)
    add_text(s, Inches(0.85), Inches(2.35), Inches(5.4), Inches(0.4),
             "RFP QUESTION", size=10, bold=True, color=MUTED)
    add_text(s, Inches(0.85), Inches(2.7), Inches(5.4), Inches(3.0),
             '"Describe your firm\'s experience delivering similar projects '
             'for public-sector clients in the last 5 years. Include outcomes '
             'and references."',
             size=14, color=BODY)

    # Right pane (AI draft)
    add_round(s, Inches(6.85), Inches(2.2), Inches(5.9), Inches(3.6),
              fill=BG, line=ORANGE)
    add_text(s, Inches(7.1), Inches(2.35), Inches(5.4), Inches(0.4),
             "AI DRAFT  ·  EDITABLE", size=10, bold=True, color=ORANGE)
    add_text(s, Inches(7.1), Inches(2.7), Inches(5.4), Inches(3.0),
             "Over the last 5 years, our firm has delivered 14 comparable "
             "engagements for public-sector clients, including the City of "
             "St. Peters water modernization (2024) and the STLCC facilities "
             "program (2023)…\n\n[Pulled from: 3 case studies, 2 references]",
             size=13, color=BODY)

    # Implement button
    btn = add_round(s, Inches(11.0), Inches(5.95), Inches(1.75), Inches(0.5),
                    fill=ORANGE)
    add_text(s, Inches(11.0), Inches(6.0), Inches(1.75), Inches(0.4),
             "▶  Implement", size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    # Chat sidecar note
    add_text(s, Inches(0.6), Inches(6.1), Inches(10), Inches(0.7),
             "💬  Chat sidecar: \"Find outcomes from our 2024 healthcare project in Missouri.\" — answer in seconds.",
             size=13, color=NAVY, bold=True)
    footer(s, num, total)
    return s


# ============================================================
# SLIDE 7 — Knowledge Flywheel
# ============================================================
def slide_kb(num, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    header_bar(s, "Build your knowledge base once. Reuse it forever.",
               kicker="The compounding asset")

    s.shapes.add_picture(str(DIAG / "03-knowledge-flywheel.png"),
                         Inches(0.6), Inches(1.7),
                         width=Inches(7.0), height=Inches(5.0))

    # Right side: what's in it
    add_text(s, Inches(8.0), Inches(1.85), Inches(5), Inches(0.4),
             "WHAT GOES IN", size=11, bold=True, color=ORANGE)
    items = [
        ("📊", "Company profile", "History, certs, bios, mission"),
        ("🛠️", "Service catalog", "Pricing, SLAs, technical specs"),
        ("📁", "Past projects", "Outcomes, testimonials, references"),
        ("📝", "Template library", "Proven \"why us\" answers"),
    ]
    y = Inches(2.35)
    for icon, title, body in items:
        add_text(s, Inches(8.0), y, Inches(0.5), Inches(0.5), icon, size=20)
        add_text(s, Inches(8.5), y, Inches(4.5), Inches(0.4),
                 title, size=15, bold=True, color=NAVY)
        add_text(s, Inches(8.5), y + Inches(0.35), Inches(4.5), Inches(0.4),
                 body, size=12, color=MUTED)
        y += Inches(0.95)

    add_text(s, Inches(8.0), Inches(6.2), Inches(5), Inches(0.6),
             "Secure & redacted. You control what gets ingested.",
             size=12, bold=True, color=BLUE)
    footer(s, num, total)
    return s


# ============================================================
# SLIDE 8 — Day in the life (before/after)
# ============================================================
def slide_day(num, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    header_bar(s, "What changes for your team.", kicker="Before / After")

    # Two columns
    add_round(s, Inches(0.6), Inches(1.85), Inches(5.95), Inches(4.9),
              fill=WHITE, line=MUTED)
    add_text(s, Inches(0.85), Inches(2.0), Inches(5.5), Inches(0.4),
             "TODAY", size=11, bold=True, color=MUTED)
    add_text(s, Inches(0.85), Inches(2.4), Inches(5.5), Inches(4.2),
             "Mon  ·  Read the 60-page RFP, build a Q checklist by hand\n\n"
             "Tue  ·  Hunt through Slack and old proposals\n\n"
             "Wed  ·  First draft, then chase SMEs for missing answers\n\n"
             "Thu  ·  Reformat to match the customer's template\n\n"
             "Fri  ·  Review, sign-off, submit — late afternoon",
             size=14, color=BODY)

    add_round(s, Inches(6.8), Inches(1.85), Inches(5.95), Inches(4.9),
              fill=NAVY, line=ORANGE)
    add_text(s, Inches(7.05), Inches(2.0), Inches(5.5), Inches(0.4),
             "WITH ROCKETRFP", size=11, bold=True, color=ORANGE)
    add_text(s, Inches(7.05), Inches(2.4), Inches(5.5), Inches(4.2),
             "Mon AM  ·  Upload RFP, review extracted checklist (15 min)\n\n"
             "Mon AM  ·  AI drafts 60–80% from your knowledge base\n\n"
             "Mon PM  ·  SMEs review the flagged gaps in the canvas\n\n"
             "Mon PM  ·  Export in original format, route for approval\n\n"
             "Tue       ·  Submitted. The rest of the week is new pipeline.",
             size=14, color=WHITE)
    footer(s, num, total)
    return s


# ============================================================
# SLIDE 9 — Workflow + Architecture
# ============================================================
def slide_workflow(num, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    header_bar(s, "End to end, in five steps.", kicker="The user flow")
    s.shapes.add_picture(str(DIAG / "02-user-workflow.png"),
                         Inches(0.6), Inches(1.7),
                         width=Inches(12.13), height=Inches(4.6))
    add_text(s, Inches(0.6), Inches(6.4), Inches(12.13), Inches(0.5),
             "Total time per response: 2–4 hours instead of 3–7 days.",
             size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    footer(s, num, total)
    return s


# ============================================================
# SLIDE 10 — Security & Multi-tenant
# ============================================================
def slide_security(num, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    header_bar(s, "Your data stays yours.", kicker="Security & isolation")
    s.shapes.add_picture(str(DIAG / "05-security-isolation.png"),
                         Inches(0.6), Inches(1.7),
                         width=Inches(8.5), height=Inches(5.0))

    add_text(s, Inches(9.4), Inches(1.85), Inches(3.7), Inches(0.4),
             "WHY IT MATTERS", size=11, bold=True, color=ORANGE)
    bullets = [
        "Per-tenant isolation — your KB never touches another customer's",
        "Encrypted at rest and in transit",
        "SSO + role-based access (admin / reviewer / contributor)",
        "Signed, expiring file URLs",
        "Optional private/local LLM for sensitive sectors",
    ]
    y = Inches(2.4)
    for b in bullets:
        add_text(s, Inches(9.4), y, Inches(0.3), Inches(0.4),
                 "▸", size=14, bold=True, color=ORANGE)
        add_text(s, Inches(9.65), y, Inches(3.4), Inches(0.9),
                 b, size=12, color=BODY)
        y += Inches(0.85)
    footer(s, num, total)
    return s


# ============================================================
# SLIDE 11 — Common questions / objections
# ============================================================
def slide_objections(num, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    header_bar(s, "What customers ask first.", kicker="Common questions")

    qa = [
        ("Does this replace my proposal team?",
         "No. It removes the copy-paste tax so they spend their time on strategy, "
         "win themes, and customer relationships — the work that actually wins bids."),
        ("How accurate is the AI-generated content?",
         "Every answer is sourced from your own knowledge base and tagged with "
         "confidence. Low-confidence items are flagged for human review before export."),
        ("What about confidentiality and our IP?",
         "Per-tenant data isolation, encryption, signed URLs, and an optional private "
         "LLM mode for regulated industries. You decide what's ingested."),
        ("How long until we see value?",
         "First completed RFP typically inside week one. The knowledge base — and the "
         "speed — compounds with every response after that."),
    ]
    x0, y0 = Inches(0.6), Inches(1.75)
    cw, ch = Inches(6.05), Inches(2.55)
    gap = Inches(0.15)
    for i, (q, a) in enumerate(qa):
        col = i % 2; row = i // 2
        x = x0 + col * (cw + gap)
        y = y0 + row * (ch + gap)
        card = add_round(s, x, y, cw, ch, fill=WHITE, line=NAVY)
        card.line.width = Pt(0.75)
        add_text(s, x + Inches(0.3), y + Inches(0.25), cw - Inches(0.6), Inches(0.85),
                 q, size=15, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), y + Inches(1.0), cw - Inches(0.6), ch - Inches(1.1),
                 a, size=12, color=BODY)
    footer(s, num, total)
    return s


# ============================================================
# SLIDE 12 — 30 days from now
# ============================================================
def slide_30days(num, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    header_bar(s, "What changes in your first 30 days.", kicker="The path to value")

    phases = [
        ("WEEK 1", "Stand up",
         "Connect your first knowledge sources. Run RocketRFP against a "
         "live RFP alongside your normal process. Compare outputs."),
        ("WEEK 2", "Switch on",
         "Your team drafts inside the canvas. SMEs review flagged gaps. "
         "First proposal submitted from RocketRFP."),
        ("WEEKS 3–4", "Scale up",
         "Multiple proposals in flight in parallel. Won outcomes feed back "
         "into the KB. Time-to-first-draft drops measurably."),
    ]
    x0, y = Inches(0.6), Inches(2.0)
    cw, ch = Inches(4.05), Inches(4.0)
    gap = Inches(0.15)
    for i, (when, title, body) in enumerate(phases):
        x = x0 + i * (cw + gap)
        card = add_round(s, x, y, cw, ch, fill=WHITE, line=ORANGE if i == 2 else NAVY)
        card.line.width = Pt(2 if i == 2 else 0.75)
        add_text(s, x + Inches(0.3), y + Inches(0.3), cw - Inches(0.6), Inches(0.4),
                 when, size=11, bold=True, color=ORANGE)
        add_text(s, x + Inches(0.3), y + Inches(0.7), cw - Inches(0.6), Inches(0.6),
                 title, size=22, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), y + Inches(1.5), cw - Inches(0.6), ch - Inches(1.7),
                 body, size=14, color=BODY)

    add_text(s, Inches(0.6), Inches(6.4), Inches(12.13), Inches(0.5),
             "Outcome: more bids answered, less senior time burned, a knowledge asset that grows.",
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    footer(s, num, total)
    return s


# ============================================================
# SLIDE 13 — CTA / Closing
# ============================================================
def slide_cta(num, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, NAVY)

    # Top accent
    add_rect(s, 0, 0, SLIDE_W, Inches(0.18), fill=ORANGE)

    wordmark(s, Inches(1.0), Inches(1.0), size=44, color=WHITE, accent=ORANGE)

    add_text(s, Inches(1.0), Inches(2.3), Inches(11), Inches(1.0),
             "Let's run RocketRFP against one of your live RFPs.",
             size=34, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(3.4), Inches(11), Inches(1.0),
             "30-minute working session. We'll show you the draft your team would have "
             "spent three days producing.",
             size=18, color=BG)

    # Next steps box
    add_round(s, Inches(1.0), Inches(4.7), Inches(11.3), Inches(1.7),
              fill=WHITE, line=ORANGE)
    add_text(s, Inches(1.3), Inches(4.85), Inches(10.7), Inches(0.4),
             "NEXT STEPS", size=11, bold=True, color=ORANGE)
    add_text(s, Inches(1.3), Inches(5.2), Inches(10.7), Inches(1.2),
             "1.  Pick one in-flight RFP to use as the test.\n"
             "2.  We connect a starter set of your knowledge base materials.\n"
             "3.  Working session — first draft on screen, end of meeting.",
             size=15, color=BODY)

    add_text(s, Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.4),
             "[ Austin's contact details — name, title, email, phone ]",
             size=12, color=ORANGE, bold=True)
    return s


# ============================================================
# Build
# ============================================================
builders = [
    slide_title,        # 1
    slide_problem,      # 2
    slide_cost,         # 3
    slide_dual_loop,    # 4
    slide_intake,       # 5
    slide_canvas,       # 6
    slide_kb,           # 7
    slide_day,          # 8
    slide_workflow,     # 9
    slide_security,     # 10
    slide_objections,   # 11
    slide_30days,       # 12
    slide_cta,          # 13
]
total = len(builders)
for i, fn in enumerate(builders, 1):
    fn(i, total)

prs.save(str(OUT))
print(f"Wrote {OUT}  ({total} slides)")
